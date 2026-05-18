import os
import sys

# Ensure python-pptx is installed
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("Installing python-pptx library...")
    os.system("pip install python-pptx")
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # Set to widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Custom Cyberpunk Theme Colors
    COLOR_BG = RGBColor(8, 7, 16)         # Deep dark background
    COLOR_TITLE = RGBColor(255, 255, 255)  # Crisp white
    COLOR_CYAN = RGBColor(0, 240, 255)    # Neon cyan
    COLOR_PURPLE = RGBColor(147, 51, 234) # Neon purple
    COLOR_MAGENTA = RGBColor(255, 0, 122) # Hot magenta
    COLOR_TEXT = RGBColor(200, 200, 215)   # Light secondary text
    COLOR_MUTED = RGBColor(120, 120, 135)  # Muted grey
    
    blank_layout = prs.slide_layouts[6]
    
    # Slides Data
    slides_data = [
        # Slide 1: Cover
        {
            "badge": "BLOCKCHAIN TECHNOLOGIES 2",
            "title": "DSA DeFi Super-App",
            "subtitle": "Full-Stack Institutional Protocol Core",
            "desc": "An institutional-grade DeFi hub combining a constant-product AMM, collateralized Lending Pool, ERC-4626 Vault, Chainlink Oracles, and DAO governance indexed via The Graph on Layer-2 networks.",
            "footer": "Team Capstone Engineering  •  L2 Verified Sandbox  •  Academic Defense Deck",
            "is_cover": True
        },
        # Slide 2: Project Overview
        {
            "badge": "01 / SCENARIO A OVERVIEW",
            "title": "Modular Web3 Ecosystem",
            "desc": "A production-ready platform linking three financial primitives under a single decentralized governance matrix, optimized for gas efficiency.",
            "points": [
                ("AMM Swap Core", "Constant product exchange (x * y = k) with slippage protections and 0.3% LP fees."),
                ("Lending Pool Engine", "Collateralized debt positions using dynamic valuations and secure liquidation gates."),
                ("ERC-4626 Yield Vault", "Standardized yield-bearing vault with automated compounding and inflation protections."),
                ("Infrastructure Stack", "Chainlink Price Feeds registry, The Graph custom indexers, and UUPS upgrades.")
            ]
        },
        # Slide 3: System Architecture
        {
            "badge": "02 / SYSTEM ARCHITECTURE",
            "title": "Modular Architecture & Data Sync Flow",
            "desc": "Our system architecture enforces separation of concerns, decoupling math algorithms, price queries, and asset storage into discrete modules.",
            "points": [
                ("UUPS Proxy Patterns", "State-changing implementations are decoupled from secure proxy storage contracts."),
                ("Separation of Concerns", "Separates pool logic, oracle feed registries, and the DAO timelock controllers."),
                ("The Graph Sync Loop", "Solidity events emit data -> Indexers parse to DB -> Frontend performs sub-10ms GQL queries."),
                ("Granular RBAC Policies", "Access control matrices secured via OpenZeppelin Role-Based permissions.")
            ]
        },
        # Slide 4: AMM Module
        {
            "badge": "03 / LIQUIDITY & AMM ENGINE",
            "title": "Constant Product Swap Core",
            "desc": "Built from scratch to deliver slippage-guarded, trustless token pair conversions matching Uniswap's core algorithm with deterministic CREATE2 deployments.",
            "points": [
                ("0.30% Liquidity Fee", "Re-invested directly into reserves to incentivize liquidity providers."),
                ("Slippage Protection Guards", "Strict assertions on minimum out parameters to prevent MEV searcher frontrunning."),
                ("Deterministic CREATE2", "AMMFactory deploys standard liquidity pairs with pre-calculated on-chain addresses."),
                ("Fuzz Tested Invariant", "Constant product invariant (k never decreases) verified through 10,000+ fuzzing runs.")
            ]
        },
        # Slide 5: Lending Pool
        {
            "badge": "04 / DEBT & COLLATERAL POOL",
            "title": "Over-Collateralized Debt System",
            "desc": "Allows users to borrow ERC-20 assets against collateral safely. Leverages dynamic valuations provided securely by our Chainlink integration.",
            "points": [
                ("75% Max LTV Threshold", "Borrow limits strictly enforced at the transaction execution boundary."),
                ("Health Factor (HF) Guards", "Reverts borrow operations instantly if account Health Factor falls below 1.0."),
                ("ERC-721 NFT Positions", "Debt profiles and collateral balances mapped directly to secure NFT assets."),
                ("Liquidation Incentives", "5% bonus awarded to third-party liquidators who pay off bad debt to protect the pool.")
            ]
        },
        # Slide 6: ERC-4626 Vault
        {
            "badge": "05 / TOKENIZED YIELD",
            "title": "Optimized Compounding Vault",
            "desc": "Standardized yield-bearing vault allowing users to deposit underlying assets (e.g. USDC) in exchange for yield-bearing shares.",
            "points": [
                ("ERC-4626 Compliance", "Fully compliant tokenized asset balances ensuring structural ecosystem interoperability."),
                ("Inflation Attack Defense", "Enforces virtual assets & shares to secure initial deposits against early depositor locks."),
                ("Automated Compounder", "Profits automatically compound inside the secure execution boundary to maximize yields."),
                ("Strict Math Invariants", "Enforces strict rounding-down criteria to prevent share extraction leakage.")
            ]
        },
        # Slide 7: Governance System
        {
            "badge": "06 / DAO GOVERNANCE",
            "title": "Autonomous Timelock Pipeline",
            "desc": "Enforces a 4-step governance state machine to secure administrative parameters without reliance on a single centralized multisig keyset.",
            "points": [
                ("ERC20Votes + Permit", "Gasless vote delegations and on-chain vote weight calculations over check snapshots."),
                ("2-Day Delay Timelock", "Allows users to withdraw assets safely if they disagree with queued parameters."),
                ("4% Quorum & 1% Threshold", "Defends against proposal spam and snap-voting flash loan takeovers."),
                ("Full Lifecycle Demonstrated", "Complete Propose -> Vote -> Queue -> Execute path fully verified.")
            ]
        },
        # Slide 8: Security & Audit
        {
            "badge": "07 / PROTOCOL SECURITY AUDIT",
            "title": "Production-Grade Safety Gates",
            "desc": "Our codebase has undergone exhaustive testing using static analyzers and manual reentrancy/access-control simulations.",
            "points": [
                ("Slither Audit clean", "0 High severity findings, 0 Medium severity findings detected across all smart contracts."),
                ("Checks-Effects-Interactions", "State transitions strictly complete before external call triggers to block attacks."),
                ("Exploit Case Studies", "Simulated Reentrancy and Access Control bypasses successfully reproduced and patched."),
                ("SafeERC20 Integrations", "All ERC-20 interactions utilize SafeERC20 to handle non-standard token transfers.")
            ]
        },
        # Slide 9: Testing Infrastructure
        {
            "badge": "08 / TESTING INFRASTRUCTURE",
            "title": "Exhaustive Multi-Layer Test Matrix",
            "desc": "High-fidelity test suite spanning four test layers. Verified fully in our GitHub Actions CI/CD pipeline.",
            "points": [
                ("50+ Unit Tests", "Complete coverage of standard execution and revert paths for all functions."),
                ("10+ Fuzz Tests", "Fuzzing input parameters for AMM swaps and Vault deposits to find edge-case reverts."),
                ("5+ Invariant Tests", "Proving mathematical invariants (e.g. constant product holds) across state transitions."),
                ("3+ Fork Tests", "Live mainnet fork testing of Chainlink feed structures and price updates.")
            ]
        },
        # Slide 10: Layer-2 Deployment
        {
            "badge": "09 / L2 SCALING & GAS MATH",
            "title": "Gas-Efficient L2 Deployments",
            "desc": "Our codebase is optimized and deployed on Base Sepolia, Arbitrum Sepolia, and Optimism Sepolia L2 testnets.",
            "points": [
                ("Yul Inline Assembly", "Math square-roots and division methods save 17-20% execution gas overhead."),
                ("Storage Slot Packing", "Packed Solidity variables minimize cold-write storage overhead slot reads."),
                ("99% Cost Reductions", "Transactions cost fractions of a cent compared to Ethereum Mainnet gas fee charts."),
                ("Idempotent Deployment", "Deployments are fully parameterized via automated deployment scripts.")
            ]
        },
        # Slide 11: Frontend & Subgraph
        {
            "badge": "10 / FRONTEND & GRAPH SYNC",
            "title": "Premium Web3 Interface & Indexers",
            "desc": "A premium responsive React dashboard built on Wagmi v2 & Apollo Client. Reads directly from custom high-performance subgraph nodes.",
            "points": [
                ("MetaMask Wallet Connect", "Automatic network checks and prompt-to-switch Sepolia network."),
                ("Apollo GraphQL Client", "Consumes indexed subgraphs offloading expensive RPC polling networks."),
                ("Custom Indexer Charts", "Displays real-time synchronization stats, query logs, and transaction streams."),
                ("Responsive Glassmorphism", "Vibrant, premium user experience matching modern institutional portals.")
            ]
        },
        # Slide 12: Conclusion
        {
            "badge": "CAPSTONE COMPLETED",
            "title": "DSA DeFi Capstone Protocol Core",
            "subtitle": "Secure, Scalable, Production-Ready",
            "desc": "We have designed, tested, audited, and deployed a production-ready Web3 platform that complies with 100% of the course requirements and holds zero high/medium security issues.",
            "footer": "Thank you! Open for Q&A session  •  Capstone Engineering 2026",
            "is_cover": True
        }
    ]
    
    for i, data in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_layout)
        
        # Add Solid Dark Background
        bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = COLOR_BG
        bg_shape.line.fill.background()
        
        if data.get("is_cover", False):
            # Layout for Cover Slides (Slide 1 & 12)
            
            # Glowing logo block (Top)
            logo_box = slide.shapes.add_textbox(Inches(5.666), Inches(1.2), Inches(2), Inches(1.2))
            tf = logo_box.text_frame
            p = tf.paragraphs[0]
            p.text = "DSA"
            p.font.name = "Space Grotesk"
            p.font.size = Pt(44)
            p.font.bold = True
            p.font.color.rgb = COLOR_CYAN
            p.alignment = PP_ALIGN.CENTER
            
            # Main Title (Center)
            title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.6), Inches(11.333), Inches(1.5))
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = data["title"]
            p.font.name = "Space Grotesk"
            p.font.size = Pt(48)
            p.font.bold = True
            p.font.color.rgb = COLOR_TITLE
            p.alignment = PP_ALIGN.CENTER
            
            # Subtitle (Below Title)
            sub_box = slide.shapes.add_textbox(Inches(1.0), Inches(3.9), Inches(11.333), Inches(0.8))
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = data["subtitle"]
            p.font.name = "Outfit"
            p.font.size = Pt(20)
            p.font.color.rgb = COLOR_PURPLE
            p.alignment = PP_ALIGN.CENTER
            
            # Description (Below Subtitle)
            desc_box = slide.shapes.add_textbox(Inches(2.0), Inches(4.7), Inches(9.333), Inches(1.5))
            tf = desc_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = data["desc"]
            p.font.name = "Outfit"
            p.font.size = Pt(14)
            p.font.color.rgb = COLOR_TEXT
            p.alignment = PP_ALIGN.CENTER
            
            # Footer (Bottom)
            foot_box = slide.shapes.add_textbox(Inches(1.0), Inches(6.5), Inches(11.333), Inches(0.6))
            tf = foot_box.text_frame
            p = tf.paragraphs[0]
            p.text = data["footer"]
            p.font.name = "Outfit"
            p.font.size = Pt(12)
            p.font.color.rgb = COLOR_MUTED
            p.alignment = PP_ALIGN.CENTER
            
        else:
            # Layout for Content Slides
            
            # Badge (Top Left)
            badge_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8), Inches(0.5))
            tf = badge_box.text_frame
            p = tf.paragraphs[0]
            p.text = data["badge"]
            p.font.name = "Space Grotesk"
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = COLOR_CYAN
            
            # Main Slide Title (Top Left)
            title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11), Inches(0.8))
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = data["title"]
            p.font.name = "Space Grotesk"
            p.font.size = Pt(26)
            p.font.bold = True
            p.font.color.rgb = COLOR_TITLE
            
            # Slide Description (Left Column)
            desc_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.2), Inches(4.5))
            tf = desc_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = data["desc"]
            p.font.name = "Outfit"
            p.font.size = Pt(15)
            p.font.color.rgb = COLOR_TEXT
            
            # Add a glowing accent bar on the left column block
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(3.6), Inches(5.2), Inches(0.04))
            bar.fill.solid()
            bar.fill.fore_color.rgb = COLOR_PURPLE
            bar.line.fill.background()
            
            # Summary / Focus Box (Lower Left)
            focus_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.9), Inches(5.2), Inches(2.2))
            tf = focus_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = "Technical Compliance Pillar:"
            p.font.name = "Space Grotesk"
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = COLOR_CYAN
            
            p2 = tf.add_paragraph()
            p2.text = "Designed and tested to fulfill 100% of the Blockchain Technologies course directives for enterprise-level decentralized architectures."
            p2.font.name = "Outfit"
            p2.font.size = Pt(12)
            p2.font.color.rgb = COLOR_MUTED
            
            # Bullet/Points Block (Right Column - styled as a beautiful technical list)
            points_box = slide.shapes.add_textbox(Inches(6.6), Inches(1.8), Inches(5.8), Inches(5))
            tf = points_box.text_frame
            tf.word_wrap = True
            
            for idx, (head, body) in enumerate(data.get("points", [])):
                p_head = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                p_head.text = f"•  {head}"
                p_head.font.name = "Space Grotesk"
                p_head.font.size = Pt(14)
                p_head.font.bold = True
                p_head.font.color.rgb = COLOR_CYAN
                p_head.space_after = Pt(2)
                if idx > 0:
                    p_head.space_before = Pt(12)
                
                p_body = tf.add_paragraph()
                p_body.text = body
                p_body.font.name = "Outfit"
                p_body.font.size = Pt(11.5)
                p_body.font.color.rgb = COLOR_TEXT
                p_body.space_after = Pt(2)
                
            # Slide Number Footer (Bottom Right)
            num_box = slide.shapes.add_textbox(Inches(11.8), Inches(6.8), Inches(1), Inches(0.4))
            tf = num_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"{i+1} / 12"
            p.font.name = "Space Grotesk"
            p.font.size = Pt(10)
            p.font.color.rgb = COLOR_MUTED
            p.alignment = PP_ALIGN.RIGHT

    # Save to file
    output_path = "/Users/assylkhanaytzhanmail.ru/Desktop/finalprojectblockchain/docs/presentation.pptx"
    prs.save(output_path)
    print(f"Presentation successfully created at: {output_path}")

if __name__ == "__main__":
    create_presentation()
